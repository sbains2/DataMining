from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import copy

# ── color palette ──────────────────────────────────────────────────────────────
DARK_BLUE  = RGBColor(0x1F, 0x45, 0x7C)   # deep navy
MID_BLUE   = RGBColor(0x2E, 0x75, 0xB6)   # title accent
LIGHT_BLUE = RGBColor(0xBD, 0xD7, 0xEE)   # subtle fill
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
DARK_GRAY  = RGBColor(0x40, 0x40, 0x40)
GREEN      = RGBColor(0x37, 0x86, 0x44)
RED        = RGBColor(0xC0, 0x00, 0x00)

prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)

BLANK = prs.slide_layouts[6]   # truly blank layout


# ── helpers ────────────────────────────────────────────────────────────────────

def add_rect(slide, l, t, w, h, fill_rgb=None, line_rgb=None, line_width_pt=0):
    shape = slide.shapes.add_shape(1, Inches(l), Inches(t), Inches(w), Inches(h))
    shape.line.fill.background()
    if fill_rgb:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_rgb
    else:
        shape.fill.background()
    if line_rgb and line_width_pt:
        from pptx.util import Pt as UPt
        shape.line.color.rgb = line_rgb
        shape.line.width = UPt(line_width_pt)
    else:
        shape.line.fill.background()
    return shape


def add_text(slide, text, l, t, w, h,
             font_size=18, bold=False, color=DARK_GRAY,
             align=PP_ALIGN.LEFT, italic=False, wrap=True):
    txb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    txb.word_wrap = wrap
    tf = txb.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return txb


def header_bar(slide, title_text, subtitle_text=""):
    add_rect(slide, 0, 0, 13.33, 1.15, fill_rgb=DARK_BLUE)
    add_text(slide, title_text, 0.35, 0.08, 12.6, 0.65,
             font_size=28, bold=True, color=WHITE, align=PP_ALIGN.LEFT)
    if subtitle_text:
        add_text(slide, subtitle_text, 0.35, 0.72, 12.6, 0.38,
                 font_size=14, color=LIGHT_BLUE, align=PP_ALIGN.LEFT)
    add_rect(slide, 0, 1.15, 13.33, 0.06, fill_rgb=MID_BLUE)


def footer(slide, page_num, total=15):
    add_rect(slide, 0, 7.15, 13.33, 0.35, fill_rgb=DARK_BLUE)
    add_text(slide, "Predicting Car Crash Severity  |  Data Mining  |  March 2026",
             0.3, 7.16, 10, 0.28, font_size=9, color=LIGHT_BLUE)
    add_text(slide, f"{page_num} / {total}", 12.5, 7.16, 0.7, 0.28,
             font_size=9, color=WHITE, align=PP_ALIGN.RIGHT)


def bullet_block(slide, items, l, t, w, h, font_size=16,
                 bullet_char="▸ ", color=DARK_GRAY, line_spacing_pt=None):
    txb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    txb.word_wrap = True
    tf = txb.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        if line_spacing_pt:
            from pptx.util import Pt as UPt
            p.line_spacing = UPt(line_spacing_pt)
        run = p.add_run()
        run.text = bullet_char + item
        run.font.size = Pt(font_size)
        run.font.color.rgb = color


def metric_box(slide, label, value, l, t, w=2.5, h=1.3,
               val_color=MID_BLUE):
    add_rect(slide, l, t, w, h, fill_rgb=LIGHT_BLUE)
    add_rect(slide, l, t, w, 0.4, fill_rgb=MID_BLUE)
    add_text(slide, label, l+0.08, t+0.04, w-0.16, 0.32,
             font_size=12, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text(slide, value, l+0.08, t+0.42, w-0.16, 0.75,
             font_size=26, bold=True, color=val_color, align=PP_ALIGN.CENTER)


def section_divider(slide, section_num, section_title, subtitle=""):
    add_rect(slide, 0, 0, 13.33, 7.5, fill_rgb=DARK_BLUE)
    add_rect(slide, 0, 2.8, 13.33, 0.08, fill_rgb=MID_BLUE)
    add_text(slide, f"Section {section_num}", 1, 1.6, 11, 0.6,
             font_size=18, color=LIGHT_BLUE, align=PP_ALIGN.CENTER)
    add_text(slide, section_title, 1, 2.2, 11, 1.0,
             font_size=42, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    if subtitle:
        add_text(slide, subtitle, 1.5, 3.3, 10, 0.7,
                 font_size=20, color=LIGHT_BLUE, align=PP_ALIGN.CENTER, italic=True)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 1 — Title
# ══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)
add_rect(slide, 0, 0, 13.33, 7.5, fill_rgb=DARK_BLUE)
add_rect(slide, 0, 3.5, 13.33, 0.08, fill_rgb=MID_BLUE)
add_rect(slide, 0, 3.58, 13.33, 3.92, fill_rgb=RGBColor(0x16, 0x32, 0x5C))

# main title
add_text(slide, "Predicting Car Crash Severity",
         0.6, 0.6, 12, 1.0, font_size=40, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
add_text(slide, "A Data Mining Analysis of California Traffic Incidents",
         0.6, 1.6, 12, 0.7, font_size=24, color=LIGHT_BLUE, align=PP_ALIGN.CENTER, italic=True)

# metadata
add_text(slide, "Course: Data Mining  |  Instructor: Dr. Leida Chen",
         1, 3.8, 11.33, 0.45, font_size=14, color=LIGHT_BLUE, align=PP_ALIGN.CENTER)
add_text(slide, "Keaton Kaess   ·   Tyler Wolf Williams   ·   Sahil Bains   ·   Will Doran   ·   Alex Piavis",
         0.6, 4.4, 12.1, 0.45, font_size=16, color=WHITE, align=PP_ALIGN.CENTER)
add_text(slide, "March 13, 2026",
         1, 5.0, 11.33, 0.4, font_size=14, color=LIGHT_BLUE, align=PP_ALIGN.CENTER)

# bottom bar
add_rect(slide, 0, 7.15, 13.33, 0.35, fill_rgb=RGBColor(0x0D, 0x1F, 0x3C))


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 2 — Agenda
# ══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)
header_bar(slide, "Agenda")
footer(slide, 2)

sections = [
    ("01", "Business Understanding",   "Why does this analysis matter?"),
    ("02", "Data Understanding",        "112,660 California crash records"),
    ("03", "Data Preparation",          "From 11 features to 98 binary columns"),
    ("04", "Modeling",                  "Decision Tree · Logistic Regression · Naive Bayes · KNN · K-Means"),
    ("05", "Evaluation & Results",      "Model comparison & key metrics"),
    ("06", "Recommendations",           "Actionable next steps for stakeholders"),
]

col_w, col_h = 3.9, 1.55
positions = [(0.35, 1.4), (4.4, 1.4), (8.45, 1.4),
             (0.35, 3.15), (4.4, 3.15), (8.45, 3.15)]

for (num, title, desc), (lx, ty) in zip(sections, positions):
    add_rect(slide, lx, ty, col_w, col_h, fill_rgb=LIGHT_BLUE)
    add_rect(slide, lx, ty, col_w, 0.38, fill_rgb=MID_BLUE)
    add_text(slide, num, lx+0.1, ty+0.04, 0.5, 0.3,
             font_size=14, bold=True, color=WHITE)
    add_text(slide, title, lx+0.65, ty+0.04, col_w-0.75, 0.3,
             font_size=13, bold=True, color=WHITE)
    add_text(slide, desc, lx+0.12, ty+0.45, col_w-0.24, 0.9,
             font_size=12, color=DARK_GRAY)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 3 — Business Understanding
# ══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)
header_bar(slide, "Business Understanding", "Why does crash severity prediction matter?")
footer(slide, 3)

# left panel
add_rect(slide, 0.35, 1.35, 5.9, 5.65, fill_rgb=RGBColor(0xF0, 0xF5, 0xFB))
add_text(slide, "The Problem", 0.5, 1.45, 5.5, 0.45,
         font_size=16, bold=True, color=DARK_BLUE)
add_text(slide,
         "Every year, California traffic crashes result in thousands of "
         "fatalities, strain emergency services, drive up insurance costs, "
         "and damage public infrastructure.\n\n"
         "Understanding what separates a minor fender-bender from a "
         "life-threatening collision has direct implications for resource "
         "allocation and policy design.",
         0.5, 1.95, 5.6, 2.5, font_size=13, color=DARK_GRAY)

add_text(slide, "Stakeholders", 0.5, 4.55, 5.5, 0.4,
         font_size=15, bold=True, color=DARK_BLUE)
bullet_block(slide,
    ["State transportation agencies (reduce fatalities)",
     "Insurance companies (model risk)",
     "Emergency dispatchers (prioritize response)",
     "Municipal planners (safer road design)"],
    0.55, 5.0, 5.55, 1.8, font_size=12)

# right panel — 3 questions
add_text(slide, "Core Business Questions", 6.6, 1.35, 6.3, 0.45,
         font_size=16, bold=True, color=DARK_BLUE)

questions = [
    ("Q1", "What factors most strongly predict severe crash outcomes?"),
    ("Q2", "Are there distinct crash risk profiles in the data?"),
    ("Q3", "How reliably can severity be predicted from available features?"),
]
for i, (qnum, qtext) in enumerate(questions):
    ty = 1.95 + i * 1.55
    add_rect(slide, 6.6, ty, 6.3, 1.3, fill_rgb=LIGHT_BLUE)
    add_rect(slide, 6.6, ty, 1.0, 1.3, fill_rgb=MID_BLUE)
    add_text(slide, qnum, 6.65, ty+0.38, 0.9, 0.5,
             font_size=18, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text(slide, qtext, 7.7, ty+0.25, 5.1, 0.9,
             font_size=13, color=DARK_GRAY)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 4 — Data Understanding
# ══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)
header_bar(slide, "Data Understanding", "112,660 California crash records · 11 original variables")
footer(slide, 4)

# stat boxes
stats = [
    ("Records", "112,660"), ("Features", "11 original"), ("Target", "Severity"),
    ("Minor (0)", "93%  ·  104,474"), ("Severe (1)", "7%  ·  8,186"),
]
for i, (lbl, val) in enumerate(stats):
    metric_box(slide, lbl, val, 0.35 + i * 2.6, 1.35, w=2.45, h=1.15)

# variable table
headers = ["Variable", "Type", "Description"]
rows = [
    ["County",        "Categorical", "CA county of occurrence (57 counties)"],
    ["Weekday",       "Ordinal",     "Day of week (1=Mon … 7=Sun)"],
    ["Severity",      "Binary",      "Target: 0=Minor, 1=Severe"],
    ["ViolCat",       "Categorical", "Violation category (11 codes)"],
    ["ClearWeather",  "Binary",      "1 if weather was clear"],
    ["Month",         "Ordinal",     "Month of occurrence (1–12)"],
    ["CrashType",     "Categorical", "Crash type A–G"],
    ["Highway",       "Binary",      "1 if on a highway"],
    ["Daylight",      "Binary",      "1 if during daylight hours"],
]

col_widths = [2.2, 1.7, 6.5]
col_starts = [0.35, 2.6, 4.35]
row_h = 0.38
t0 = 2.7

for ci, (hdr, cw, cs) in enumerate(zip(headers, col_widths, col_starts)):
    add_rect(slide, cs, t0, cw, row_h, fill_rgb=DARK_BLUE)
    add_text(slide, hdr, cs+0.05, t0+0.04, cw-0.1, row_h-0.06,
             font_size=12, bold=True, color=WHITE)

for ri, row in enumerate(rows):
    ty = t0 + row_h + ri * row_h
    fill = LIGHT_BLUE if ri % 2 == 0 else WHITE
    for ci, (cell, cw, cs) in enumerate(zip(row, col_widths, col_starts)):
        add_rect(slide, cs, ty, cw, row_h, fill_rgb=fill)
        fc = DARK_BLUE if cell == "Severity" and ci == 0 else DARK_GRAY
        fb = cell == "Severity" and ci == 0
        add_text(slide, cell, cs+0.06, ty+0.04, cw-0.1, row_h-0.06,
                 font_size=11, color=fc, bold=fb)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 5 — Class Imbalance
# ══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)
header_bar(slide, "Class Imbalance Challenge",
           "93% minor vs 7% severe — why accuracy alone is misleading")
footer(slide, 5)

# big bar visual
add_rect(slide, 0.5, 1.45, 11.5, 0.6, fill_rgb=RGBColor(0xE0, 0xE0, 0xE0))
add_rect(slide, 0.5, 1.45, 10.695, 0.6, fill_rgb=MID_BLUE)   # 93%
add_rect(slide, 11.195, 1.45, 0.805, 0.6, fill_rgb=RED)      # 7%
add_text(slide, "Minor (93%)", 3.5, 1.52, 5, 0.45,
         font_size=16, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
add_text(slide, "Severe 7%", 11.0, 1.52, 1.2, 0.45,
         font_size=11, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

# two columns of explanation
add_text(slide, "The Trap", 0.5, 2.25, 5.9, 0.4,
         font_size=15, bold=True, color=RED)
add_text(slide,
         "A naive classifier that always predicts 'Minor' achieves 93% accuracy "
         "— yet is completely useless. It would miss every life-threatening crash.",
         0.5, 2.7, 5.9, 1.2, font_size=13, color=DARK_GRAY)

add_text(slide, "Our Approach", 7.0, 2.25, 5.9, 0.4,
         font_size=15, bold=True, color=GREEN)
add_text(slide,
         "We prioritise AUC-ROC, F1-score, and severe-class recall over raw accuracy. "
         "Minority class upsampling (Decision Tree) and balanced class priors "
         "(Naive Bayes, Logistic Regression) were applied to counteract the imbalance.",
         7.0, 2.7, 5.9, 1.6, font_size=13, color=DARK_GRAY)

# divider
add_rect(slide, 6.4, 2.25, 0.06, 4.5, fill_rgb=LIGHT_BLUE)

bullet_block(slide,
    ["Minor crashes: 104,474 records",
     "Severe crashes:    8,186 records",
     "Imbalance ratio:   ~13 : 1"],
    0.5, 4.0, 5.9, 1.6, font_size=14)

bullet_block(slide,
    ["AUC-ROC > accuracy for evaluation",
     "F1-score balances precision & recall",
     "Severe-class recall = critical metric"],
    7.0, 4.0, 5.9, 1.6, font_size=14)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 6 — Data Preparation
# ══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)
header_bar(slide, "Data Preparation", "Three-stage pipeline: clean → encode → fully binary")
footer(slide, 6)

stages = [
    ("Stage 1\nInitial Cleaning", [
        "Dropped ID (non-predictive identifier)",
        "Removed City (redundant with County)",
        "Addressed missing / null values",
        "Dataset was largely clean — no major imputation needed",
    ]),
    ("Stage 2\nPartial Encoding", [
        "One-hot encoded CrashType (A–G) → 7 columns",
        "One-hot encoded ViolCat (11 codes) → 11 columns",
        "Result: 25-column intermediate dataset",
        "Ordinal integers kept for County, Weekday, Month",
    ]),
    ("Stage 3\nFull Binary Encoding", [
        "One-hot encoded County, Weekday, Month",
        "All boolean columns cast to 0/1 integers",
        "Final dataset: 112,660 rows × 98 columns",
        "Used as input for ALL five models",
    ]),
]

arrow_color = MID_BLUE
sw = 3.8
gap = 0.3
lx = 0.35

for i, (title, bullets) in enumerate(stages):
    x = lx + i * (sw + gap)
    add_rect(slide, x, 1.4, sw, 5.5, fill_rgb=RGBColor(0xF0, 0xF5, 0xFB))
    add_rect(slide, x, 1.4, sw, 0.7, fill_rgb=MID_BLUE)
    add_text(slide, title, x+0.1, 1.42, sw-0.2, 0.65,
             font_size=13, bold=True, color=WHITE)
    bullet_block(slide, bullets, x+0.15, 2.2, sw-0.25, 4.4,
                 font_size=12, bullet_char="• ")
    if i < 2:
        add_text(slide, "▶", x+sw+0.05, 3.8, 0.28, 0.5,
                 font_size=20, bold=True, color=MID_BLUE, align=PP_ALIGN.CENTER)

# summary badge
add_rect(slide, 0.35, 6.6, 12.6, 0.55, fill_rgb=DARK_BLUE)
add_text(slide,
         "Result: 112,660 records  ·  98 fully binary features  ·  ready for all classifiers",
         0.5, 6.63, 12.3, 0.45, font_size=13, bold=True, color=WHITE, align=PP_ALIGN.CENTER)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 7 — Modeling Overview
# ══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)
header_bar(slide, "Modeling Approach", "Five algorithms across supervised and unsupervised paradigms")
footer(slide, 7)

models = [
    ("Decision\nTree",       MID_BLUE,  "Supervised",   "max_depth=6\nBalanced weights\nUpsampled minority\n5-fold CV"),
    ("Logistic\nRegression", MID_BLUE,  "Supervised",   "L-BFGS solver\nStandardScaler\nBalanced weights\n75/25 split"),
    ("Naive\nBayes",         MID_BLUE,  "Supervised",   "Bernoulli NB\nalpha=1.0 (Laplace)\nEqual class priors\nNo scaling needed"),
    ("KNN\n(k=5)",           MID_BLUE,  "Supervised",   "Euclidean distance\nDistance weighting\nStandardScaler\n70/30 split"),
    ("K-Means\n(K=3)",       GREEN,     "Unsupervised", "PCA → 10 components\n84.3% variance kept\nElbow method\nSilhouette: 0.204"),
]

bw = 2.3
for i, (name, col, paradigm, config) in enumerate(models):
    x = 0.35 + i * (bw + 0.2)
    add_rect(slide, x, 1.4, bw, 5.4, fill_rgb=RGBColor(0xF0, 0xF5, 0xFB))
    add_rect(slide, x, 1.4, bw, 0.75, fill_rgb=col)
    add_text(slide, name, x+0.08, 1.42, bw-0.16, 0.7,
             font_size=14, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    badge_col = MID_BLUE if paradigm == "Supervised" else GREEN
    add_rect(slide, x+0.15, 2.22, bw-0.3, 0.32, fill_rgb=badge_col)
    add_text(slide, paradigm, x+0.15, 2.23, bw-0.3, 0.3,
             font_size=10, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text(slide, config, x+0.12, 2.65, bw-0.2, 3.8,
             font_size=11, color=DARK_GRAY)

add_text(slide,
         "Evaluation metrics: AUC-ROC · F1-Score · Recall (Severe class) · Precision · Overall Accuracy",
         0.35, 6.6, 12.6, 0.45, font_size=12, color=DARK_GRAY, align=PP_ALIGN.CENTER, italic=True)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 8 — Decision Tree Results
# ══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)
header_bar(slide, "Decision Tree — Best Performing Model",
           "AUC-ROC: 0.838 · F1: 0.853 · Accuracy: 89.1% · 5-Fold CV F1: 0.812 ± 0.034")
footer(slide, 8)

# metric boxes
metrics = [("AUC-ROC", "0.838"), ("F1-Score", "0.853"),
           ("Accuracy", "89.1%"), ("CV F1", "0.812 ± 0.034")]
for i, (lbl, val) in enumerate(metrics):
    metric_box(slide, lbl, val, 0.35 + i * 3.2, 1.42, w=3.0, h=1.15)

# left: feature importances
add_text(slide, "Top Feature Importances", 0.35, 2.75, 5.8, 0.38,
         font_size=14, bold=True, color=DARK_BLUE)

features = [
    ("Highway",      0.1423, MID_BLUE),
    ("Daylight",     0.1187, MID_BLUE),
    ("ClearWeather", 0.0891, MID_BLUE),
    ("ViolCat_1",    0.0754, MID_BLUE),
    ("CrashType_A",  0.0632, MID_BLUE),
    ("County_28",    0.0481, RGBColor(0x70, 0xA0, 0xD0)),
    ("ViolCat_4",    0.0398, RGBColor(0x70, 0xA0, 0xD0)),
]
bar_max_w = 4.5
for j, (fname, fval, col) in enumerate(features):
    ty = 3.25 + j * 0.48
    bar_w = fval / 0.15 * bar_max_w
    add_rect(slide, 1.55, ty, bar_w, 0.34, fill_rgb=col)
    add_text(slide, fname, 0.38, ty+0.04, 1.1, 0.3, font_size=11, color=DARK_GRAY)
    add_text(slide, f"{fval:.4f}", 1.58 + bar_w, ty+0.04, 0.8, 0.3,
             font_size=10, color=DARK_GRAY)

# right: config & insights
add_text(slide, "Configuration & Insights", 6.55, 2.75, 6.4, 0.38,
         font_size=14, bold=True, color=DARK_BLUE)
bullet_block(slide,
    ["max_depth=6 to prevent overfitting",
     "min_samples_split=50, min_samples_leaf=25",
     "Balanced class weights + minority upsampling",
     "80/20 stratified train/test split",
     "",
     "Highway crashes are the single strongest predictor",
     "Lighting (Daylight) and weather (ClearWeather) follow",
     "Specific violation categories signal elevated risk",
     "Interpretable tree structure supports stakeholder audit"],
    6.55, 3.2, 6.35, 3.6, font_size=12)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 9 — Logistic Regression & Naive Bayes
# ══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)
header_bar(slide, "Logistic Regression & Naive Bayes",
           "Both deliver acceptable performance — AUC ~0.75, Recall ~65%")
footer(slide, 9)

# LR column
add_rect(slide, 0.35, 1.42, 6.0, 5.55, fill_rgb=RGBColor(0xF0, 0xF5, 0xFB))
add_rect(slide, 0.35, 1.42, 6.0, 0.5, fill_rgb=MID_BLUE)
add_text(slide, "Logistic Regression", 0.45, 1.44, 5.8, 0.45,
         font_size=15, bold=True, color=WHITE)

lr_metrics = [("AUC-ROC","0.753"), ("Accuracy","71%"), ("Recall (Sev)","66%"), ("Precision (Sev)","15%")]
for i, (lbl, val) in enumerate(lr_metrics):
    metric_box(slide, lbl, val, 0.45 + i * 1.47, 2.05, w=1.35, h=1.0)

bullet_block(slide,
    ["L-BFGS solver · StandardScaler · balanced class weights",
     "Converged in 13 iterations",
     "Of 1,979 severe crashes → 1,300 correctly identified",
     "Coefficients rank features by severity risk contribution",
     "Useful for policy targeting via coefficient magnitudes"],
    0.45, 3.2, 5.8, 2.6, font_size=12)

# NB column
add_rect(slide, 6.9, 1.42, 6.0, 5.55, fill_rgb=RGBColor(0xF0, 0xF5, 0xFB))
add_rect(slide, 6.9, 1.42, 6.0, 0.5, fill_rgb=MID_BLUE)
add_text(slide, "Naive Bayes (Bernoulli)", 7.0, 1.44, 5.8, 0.45,
         font_size=15, bold=True, color=WHITE)

nb_metrics = [("AUC-ROC","0.745"), ("Accuracy","71%"), ("Recall (Sev)","64%"), ("Precision (Sev)","14%")]
for i, (lbl, val) in enumerate(nb_metrics):
    metric_box(slide, lbl, val, 7.0 + i * 1.47, 2.05, w=1.35, h=1.0)

bullet_block(slide,
    ["Bernoulli NB — natural fit for 98 binary features",
     "Laplace smoothing alpha=1.0; class priors [0.5, 0.5]",
     "No feature scaling required",
     "TP: 1,267  ·  FN: 712  ·  TN: 18,645  ·  FP: 7,541",
     "Log-probability ratio reveals top risk features"],
    7.0, 3.2, 5.8, 2.6, font_size=12)

# divider
add_rect(slide, 6.5, 1.42, 0.08, 5.55, fill_rgb=LIGHT_BLUE)

add_text(slide, "Both models recover ~65% of severe crashes and are suitable for probability-based risk scoring.",
         0.35, 6.75, 12.6, 0.35, font_size=11, color=DARK_GRAY, align=PP_ALIGN.CENTER, italic=True)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 10 — KNN & K-Means
# ══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)
header_bar(slide, "KNN & K-Means Clustering",
           "KNN fails under class imbalance · K-Means reveals three crash archetypes")
footer(slide, 10)

# KNN
add_rect(slide, 0.35, 1.42, 5.8, 5.55, fill_rgb=RGBColor(0xF0, 0xF5, 0xFB))
add_rect(slide, 0.35, 1.42, 5.8, 0.5, fill_rgb=RED)
add_text(slide, "KNN  (k=5)  — Poor", 0.45, 1.44, 5.6, 0.45,
         font_size=15, bold=True, color=WHITE)

knn_metrics = [("AUC-ROC","0.607"), ("Accuracy","92%"), ("Recall (Sev)","6%")]
for i, (lbl, val) in enumerate(knn_metrics):
    vc = RED if lbl in ("AUC-ROC","Recall (Sev)") else GREEN
    metric_box(slide, lbl, val, 0.45 + i * 1.9, 2.05, w=1.75, h=1.0, val_color=vc)

bullet_block(slide,
    ["92% accuracy is entirely driven by majority class",
     "Only 142 of 2,375 severe crashes identified",
     "Root cause: 93/7 imbalance dominates neighborhoods",
     "McNemar test vs LR: p = 2.09 × 10⁻¹⁸⁴",
     "Not recommended in current form",
     "Future fix: SMOTE + threshold calibration"],
    0.45, 3.2, 5.65, 3.5, font_size=12)

# K-Means
add_rect(slide, 6.75, 1.42, 6.2, 5.55, fill_rgb=RGBColor(0xF0, 0xF5, 0xFB))
add_rect(slide, 6.75, 1.42, 6.2, 0.5, fill_rgb=GREEN)
add_text(slide, "K-Means Clustering (K=3)", 6.85, 1.44, 6.0, 0.45,
         font_size=15, bold=True, color=WHITE)

add_text(slide, "Silhouette Score: 0.204  ·  PCA retained 84.3% variance  ·  Elbow at K=3",
         6.85, 2.05, 5.9, 0.4, font_size=11, color=DARK_GRAY, italic=True)

clusters = [
    ("Cluster 0", "Lower-severity non-highway crashes.\nClear weather, daytime, moderate violation rate.\nSeverity: 5.2%",
     RGBColor(0x5B, 0x9B, 0xD5)),
    ("Cluster 1", "HIGHEST RISK segment.\nHighway-dominant, highest violation rate.\nSeverity: 8.9%",
     RED),
    ("Cluster 2", "Surface-street, clear-weather daytime.\nLowest highway %, moderate severity.\nSeverity: 7.1%",
     GREEN),
]
for i, (clabel, cdesc, col) in enumerate(clusters):
    ty = 2.6 + i * 1.35
    add_rect(slide, 6.85, ty, 5.95, 1.2, fill_rgb=LIGHT_BLUE)
    add_rect(slide, 6.85, ty, 1.5, 1.2, fill_rgb=col)
    add_text(slide, clabel, 6.88, ty+0.35, 1.45, 0.45,
             font_size=12, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text(slide, cdesc, 8.45, ty+0.1, 4.2, 1.0, font_size=11, color=DARK_GRAY)

add_rect(slide, 6.5, 1.42, 0.08, 5.55, fill_rgb=LIGHT_BLUE)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 11 — Model Comparison
# ══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)
header_bar(slide, "Model Performance Comparison",
           "Decision Tree dominates; KNN unsuitable without imbalance correction")
footer(slide, 11)

# bar chart — AUC-ROC
model_data = [
    ("Decision Tree",       0.838, 0.853, "High", DARK_BLUE),
    ("Logistic Regression", 0.753, 0.820, "66%",  MID_BLUE),
    ("Naive Bayes",         0.745, 0.820, "64%",  MID_BLUE),
    ("KNN (k=5)",           0.607, 0.090, "6%",   RED),
]

chart_l, chart_t, chart_h = 0.5, 1.5, 3.8
chart_w_max = 5.5

add_text(slide, "AUC-ROC Score", 0.5, 1.42, 6.0, 0.38,
         font_size=14, bold=True, color=DARK_BLUE)
add_rect(slide, 0.5, 1.8, 0.02, chart_h, fill_rgb=DARK_GRAY)  # y-axis

for i, (name, auc, f1, recall, col) in enumerate(model_data):
    ty = chart_t + 0.2 + i * 0.85
    bar_w = auc * chart_w_max
    add_rect(slide, 0.55, ty, bar_w, 0.55, fill_rgb=col)
    add_text(slide, f"{auc}", 0.55 + bar_w + 0.05, ty+0.1, 0.65, 0.35,
             font_size=12, bold=True, color=col)

# label axis
for i, (name, *_) in enumerate(model_data):
    ty = chart_t + 0.2 + i * 0.85
    pass  # names go in table below

# comparison table
add_text(slide, "Full Comparison", 0.35, 5.55, 12.6, 0.38,
         font_size=14, bold=True, color=DARK_BLUE)

t_headers = ["Model", "AUC-ROC", "Accuracy", "Recall (Severe)", "F1-Score", "Rating"]
t_col_w =   [3.2,      1.6,       1.6,        2.2,              1.6,        1.6]
t_col_x = [0.35]
for w in t_col_w[:-1]:
    t_col_x.append(t_col_x[-1] + w)

row_h2 = 0.36
t0 = 5.95
for ci, (hdr, cw, cx) in enumerate(zip(t_headers, t_col_w, t_col_x)):
    add_rect(slide, cx, t0, cw, row_h2, fill_rgb=DARK_BLUE)
    add_text(slide, hdr, cx+0.05, t0+0.04, cw-0.1, row_h2-0.06,
             font_size=11, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

t_rows = [
    ["Decision Tree",       "0.838", "89.1%", "High", "0.853", "Good"],
    ["Logistic Regression", "0.753", "71%",   "66%",  "0.820", "Acceptable"],
    ["Naive Bayes",         "0.745", "71%",   "64%",  "0.820", "Acceptable"],
    ["KNN (k=5)",           "0.607", "92%",   "6%",   "0.090", "Poor"],
]
row_colors = [
    [RGBColor(0xD9, 0xEA, 0xF7), DARK_BLUE, DARK_BLUE, DARK_BLUE, DARK_BLUE, GREEN],
    [LIGHT_BLUE]*5 + [MID_BLUE],
    [LIGHT_BLUE]*5 + [MID_BLUE],
    [RGBColor(0xFC, 0xE4, 0xD6)]*5 + [RED],
]
for ri, (row, rcolors) in enumerate(zip(t_rows, row_colors)):
    ty = t0 + row_h2 + ri * row_h2
    fill = LIGHT_BLUE if ri % 2 == 0 else WHITE
    for ci, (cell, cw, cx, rc) in enumerate(zip(row, t_col_w, t_col_x, rcolors)):
        add_rect(slide, cx, ty, cw, row_h2, fill_rgb=fill)
        fc = rc if ci in (1, 5) else DARK_GRAY
        add_text(slide, cell, cx+0.05, ty+0.03, cw-0.1, row_h2-0.05,
                 font_size=11, color=fc, bold=(ci == 5), align=PP_ALIGN.CENTER)

# right side — visual AUC bars legend with names
for i, (name, auc, f1, recall, col) in enumerate(model_data):
    ty = chart_t + 0.2 + i * 0.85
    add_text(slide, name, 6.5, ty+0.08, 3.5, 0.4, font_size=12, color=col, bold=True)
    add_text(slide, f"Recall: {recall}   F1: {f1}", 6.5, ty+0.47, 3.5, 0.3,
             font_size=10, color=DARK_GRAY)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 12 — Key Findings (Answering Business Questions)
# ══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)
header_bar(slide, "Key Findings", "Returning to our three core business questions")
footer(slide, 12)

findings = [
    ("Q1", "What predicts severe crashes?",
     ["Highway crashes are the #1 predictor (importance: 0.142)",
      "Nighttime crashes (Daylight=0) strongly elevate risk",
      "Adverse weather and specific violation categories signal severity",
      "Both Decision Tree and Logistic Regression agree on top features"]),
    ("Q2", "Are there distinct risk profiles?",
     ["K-Means found 3 natural archetypes (Silhouette = 0.204)",
      "Cluster 1: highway + high violation rate → 8.9% severity (highest)",
      "Cluster 0: non-highway, clear, daytime → 5.2% severity (lowest)",
      "Profiles enable tailored interventions, not one-size-fits-all"]),
    ("Q3", "How reliable is the prediction?",
     ["Decision Tree AUC-ROC = 0.838 — solid real-world performance",
      "~12–15% of severe crashes still missed (inherent data ceiling)",
      "Far exceeds random chance (AUC = 0.50)",
      "Sufficient for meaningful risk triage and dispatch support"]),
]

for i, (qnum, qtitle, qbullets) in enumerate(findings):
    lx = 0.35 + i * 4.35
    add_rect(slide, lx, 1.45, 4.15, 5.55, fill_rgb=RGBColor(0xF0, 0xF5, 0xFB))
    add_rect(slide, lx, 1.45, 4.15, 0.65, fill_rgb=DARK_BLUE)
    add_text(slide, qnum, lx+0.08, 1.47, 0.6, 0.6,
             font_size=20, bold=True, color=LIGHT_BLUE)
    add_text(slide, qtitle, lx+0.72, 1.52, 3.35, 0.5,
             font_size=13, bold=True, color=WHITE)
    bullet_block(slide, qbullets, lx+0.15, 2.2, 3.9, 4.5,
                 font_size=12, bullet_char="✓ ", color=DARK_GRAY)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 13 — Recommendations
# ══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)
header_bar(slide, "Recommendations", "Four actionable steps for traffic safety stakeholders")
footer(slide, 13)

recs = [
    ("1", "Deploy Decision Tree as Real-Time Triage",
     MID_BLUE,
     "Integrate the Decision Tree into emergency dispatch or accident reporting systems. "
     "Flag likely-severe crashes in real time based on road type, lighting, and violation "
     "category. AUC-ROC of 0.838 provides meaningful risk differentiation."),
    ("2", "Use Logistic Regression for Policy Targeting",
     MID_BLUE,
     "Coefficient magnitudes directly quantify each feature's marginal contribution to "
     "severity risk. Features with the largest positive coefficients should anchor "
     "enforcement campaigns, driver education programs, and infrastructure investments."),
    ("3", "Design Cluster-Specific Prevention Programs",
     GREEN,
     "The three K-Means archetypes carry distinct risk signatures. Cluster 1 (highway + "
     "high violation) requires targeted highway patrol. Cluster 0 warrants different "
     "countermeasures. Tailored programs outperform uniform interventions."),
    ("4", "Address Class Imbalance in Future Iterations",
     RGBColor(0xD0, 0x60, 0x00),
     "All classifiers were challenged by the 93/7 split. Future work should explore SMOTE "
     "oversampling and decision-threshold calibration. Ensemble methods (Random Forest, "
     "Gradient Boosting) may also significantly improve severe-class recall."),
]

for i, (num, title, col, desc) in enumerate(recs):
    ty = 1.45 + i * 1.38
    add_rect(slide, 0.35, ty, 12.6, 1.25, fill_rgb=RGBColor(0xF0, 0xF5, 0xFB))
    add_rect(slide, 0.35, ty, 0.7, 1.25, fill_rgb=col)
    add_text(slide, num, 0.38, ty+0.35, 0.62, 0.5,
             font_size=22, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text(slide, title, 1.15, ty+0.1, 4.5, 0.4,
             font_size=13, bold=True, color=col)
    add_text(slide, desc, 1.15, ty+0.52, 11.65, 0.65,
             font_size=12, color=DARK_GRAY)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 14 — Limitations & Future Work
# ══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)
header_bar(slide, "Limitations & Future Work", "Where we are and where to go next")
footer(slide, 14)

add_text(slide, "Current Limitations", 0.35, 1.45, 6.0, 0.4,
         font_size=15, bold=True, color=RED)
bullet_block(slide,
    ["93/7 class imbalance limits all classifiers, especially KNN",
     "Binary encoding discards within-county geographic nuance",
     "No continuous features (speed, age, BAC) available",
     "Dataset is historical — real-time deployment requires retraining",
     "KNN infeasible at scale without dimensionality reduction",
     "Silhouette score of 0.204 indicates soft cluster boundaries"],
    0.35, 1.95, 5.9, 4.5, font_size=13)

add_rect(slide, 6.4, 1.45, 0.06, 5.3, fill_rgb=LIGHT_BLUE)

add_text(slide, "Future Work", 6.75, 1.45, 6.2, 0.4,
         font_size=15, bold=True, color=GREEN)
bullet_block(slide,
    ["SMOTE oversampling + decision-threshold calibration",
     "Ensemble methods: Random Forest, Gradient Boosting (XGBoost)",
     "Add continuous features: speed limit, driver age, road condition",
     "Geospatial clustering to replace county-level encoding",
     "Online learning pipeline for real-time severity scoring",
     "Explainability layer (SHAP values) for dispatch tool trust"],
    6.75, 1.95, 6.2, 4.5, font_size=13)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 15 — Thank You / Q&A
# ══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)
add_rect(slide, 0, 0, 13.33, 7.5, fill_rgb=DARK_BLUE)
add_rect(slide, 0, 3.2, 13.33, 0.08, fill_rgb=MID_BLUE)

add_text(slide, "Thank You", 0.6, 0.5, 12, 1.1,
         font_size=52, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
add_text(slide, "Questions?", 0.6, 1.65, 12, 0.7,
         font_size=28, color=LIGHT_BLUE, align=PP_ALIGN.CENTER, italic=True)

# summary strip
add_rect(slide, 0.8, 3.4, 11.73, 2.6, fill_rgb=RGBColor(0x16, 0x32, 0x5C))
summary_items = [
    "Dataset: 112,660 California crash records",
    "Best model: Decision Tree  (AUC 0.838, F1 0.853)",
    "Key predictors: Highway · Daylight · ClearWeather · ViolCat",
    "3 crash risk archetypes identified via K-Means",
    "Crash severity IS predictable — ready for real-world deployment",
]
bullet_block(slide, summary_items, 1.0, 3.55, 11.3, 2.3,
             font_size=13, color=LIGHT_BLUE, bullet_char="  ◆  ")

add_text(slide,
         "Keaton Kaess  ·  Tyler Wolf Williams  ·  Sahil Bains  ·  Will Doran  ·  Alex Piavis",
         0.6, 6.2, 12.1, 0.4, font_size=13, color=LIGHT_BLUE, align=PP_ALIGN.CENTER)
add_text(slide, "Data Mining  ·  Dr. Leida Chen  ·  March 2026",
         0.6, 6.65, 12.1, 0.4, font_size=11, color=RGBColor(0x80, 0xA0, 0xC0),
         align=PP_ALIGN.CENTER)


# ── save ───────────────────────────────────────────────────────────────────────
out_path = r"c:\Users\tyler\OneDrive\Desktop\MSBA\Data Mining and Analytics\Final Project\DataMining\CarCrashSeverity_Presentation.pptx"
prs.save(out_path)
print(f"Saved: {out_path}")
